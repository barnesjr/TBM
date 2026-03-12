import json
from pathlib import Path
from datetime import datetime
from typing import Optional


MATURITY_BANDS = [
    {"min": 1.0, "max": 1.75, "label": "Ad Hoc", "color": "#ef4444"},
    {"min": 1.75, "max": 2.5, "label": "Foundational", "color": "#f97316"},
    {"min": 2.5, "max": 3.25, "label": "Managed", "color": "#84cc16"},
    {"min": 3.25, "max": 4.0, "label": "Optimized", "color": "#22c55e"},
]

SCORE_LABELS = {1: "Ad Hoc", 2: "Foundational", 3: "Managed", 4: "Optimized"}


class ExportEngine:
    def __init__(self, base_dir: str, resource_dir: str | None = None):
        self.base_dir = Path(base_dir)
        self.resource_dir = Path(resource_dir) if resource_dir else self.base_dir
        self.exports_dir = self.base_dir / "exports"
        self.templates_dir = self.resource_dir / "templates"

    def _ensure_exports_dir(self):
        self.exports_dir.mkdir(exist_ok=True)

    def _score_avg(self, items: list) -> Optional[float]:
        scored = [i for i in items if i.get("score") is not None and not i.get("na", False)]
        if not scored:
            return None
        return sum(i["score"] for i in scored) / len(scored)

    def _get_maturity_band(self, score: float) -> dict:
        for band in MATURITY_BANDS:
            if band["min"] <= score < band["max"]:
                return band
        if score >= 4.0:
            return MATURITY_BANDS[-1]
        return MATURITY_BANDS[0]

    def _weighted_composite(self, data: dict) -> Optional[float]:
        weights = data.get("scoring_config", {}).get("discipline_weights", {})
        total_weight = 0.0
        weighted_sum = 0.0
        for disc in data.get("disciplines", []):
            if not disc.get("enabled", True):
                continue
            items = []
            for ca in disc.get("capability_areas", []):
                items.extend(ca.get("items", []))
            score = self._score_avg(items)
            weight = weights.get(disc["id"], 0)
            if score is not None:
                weighted_sum += score * weight
                total_weight += weight
        if total_weight == 0:
            return None
        return weighted_sum / total_weight

    def _timestamp(self) -> str:
        return datetime.now().strftime("%Y-%m-%d_%H%M%S")

    def _enabled_disciplines(self, data: dict) -> list:
        return [d for d in data.get("disciplines", []) if d.get("enabled", True)]

    def generate_radar_chart_png(self, data: dict) -> str:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        import numpy as np

        labels = []
        scores = []
        targets = []
        target_scores = data.get("target_scores", {})
        for disc in self._enabled_disciplines(data):
            labels.append(disc["name"])
            items = []
            for ca in disc.get("capability_areas", []):
                items.extend(ca.get("items", []))
            s = self._score_avg(items)
            scores.append(s if s is not None else 0)
            targets.append(target_scores.get(disc["id"], 3.0))

        n = len(labels)
        if n == 0:
            fig, ax = plt.subplots(figsize=(6, 6))
            ax.text(0.5, 0.5, "No data", ha="center", va="center")
            self._ensure_exports_dir()
            path = str(self.exports_dir / "radar_chart.png")
            fig.savefig(path, dpi=150, bbox_inches="tight", facecolor="white")
            plt.close(fig)
            return path

        angles = np.linspace(0, 2 * np.pi, n, endpoint=False).tolist()
        scores_plot = scores + [scores[0]]
        targets_plot = targets + [targets[0]]
        angles_plot = angles + [angles[0]]

        fig, ax = plt.subplots(figsize=(6, 6), subplot_kw=dict(polar=True))
        ax.fill(angles_plot, scores_plot, alpha=0.2, color="#1BA1E2")
        ax.plot(angles_plot, scores_plot, color="#1BA1E2", linewidth=2, label="Current")
        ax.plot(angles_plot, targets_plot, color="#8A8A8E", linewidth=1.5, linestyle="--", label="Target")
        ax.set_xticks(angles)
        ax.set_xticklabels(labels, size=7)
        ax.set_ylim(0, 4)
        ax.set_yticks([1, 2, 3, 4])
        ax.set_yticklabels(["1", "2", "3", "4"], size=8)
        ax.set_title("TBM Maturity Profile", size=14, pad=20)
        ax.legend(loc="upper right", bbox_to_anchor=(1.3, 1.1), fontsize=9)

        self._ensure_exports_dir()
        path = str(self.exports_dir / "radar_chart.png")
        fig.savefig(path, dpi=150, bbox_inches="tight", facecolor="white")
        plt.close(fig)
        return path

    def export_findings(self, data: dict) -> str:
        self._ensure_exports_dir()
        filename = f"D-01_TBM_Findings_{self._timestamp()}.docx"
        output_path = self.exports_dir / filename

        from docx import Document
        from docx.shared import Pt
        doc = Document()
        doc.add_heading("Technology Business Management Assessment", 0)
        doc.add_heading("Detailed Findings", level=1)

        info = data.get("client_info", {})
        doc.add_paragraph(f"Client: {info.get('name', '')}")
        doc.add_paragraph(f"Industry: {info.get('industry', '')}")
        doc.add_paragraph(f"Date: {info.get('assessment_date', '')}")
        doc.add_paragraph(f"Assessor: {info.get('assessor', '')}")

        composite = self._weighted_composite(data)
        if composite:
            band = self._get_maturity_band(composite)
            doc.add_heading("Overall Maturity", level=1)
            doc.add_paragraph(f"Composite Score: {composite:.2f} — {band['label']}")

        for disc in self._enabled_disciplines(data):
            items_all = []
            for ca in disc.get("capability_areas", []):
                items_all.extend(ca.get("items", []))
            score = self._score_avg(items_all)
            doc.add_heading(disc["name"], level=1)
            doc.add_paragraph(
                f"Score: {score:.2f} — {self._get_maturity_band(score)['label']}" if score else "Score: Not yet scored"
            )

            for ca in disc.get("capability_areas", []):
                ca_score = self._score_avg(ca.get("items", []))
                doc.add_heading(ca["name"], level=2)
                doc.add_paragraph(f"Average Score: {ca_score:.2f}" if ca_score else "Not scored")

                for item in ca.get("items", []):
                    score_val = item.get("score")
                    na = item.get("na", False)
                    if na:
                        text = f"[N/A] {item['text']}"
                    elif score_val:
                        text = f"[{score_val} - {SCORE_LABELS.get(score_val, '')}] {item['text']}"
                    else:
                        text = f"[--] {item['text']}"
                    if item.get("notes"):
                        text += f"\n  Notes: {item['notes']}"
                    refs = item.get("evidence_references", [])
                    if refs:
                        evidence_str = "; ".join(
                            f"{r.get('document', '')} §{r.get('section', '')}"
                            for r in refs if r.get("document")
                        )
                        if evidence_str:
                            text += f"\n  Evidence: {evidence_str}"
                    doc.add_paragraph(text, style="List Bullet")

        doc.save(str(output_path))
        return filename

    def export_executive_summary(self, data: dict) -> str:
        self._ensure_exports_dir()
        filename = f"D-02_TBM_Executive_Summary_{self._timestamp()}.docx"
        output_path = self.exports_dir / filename
        chart_path = self.generate_radar_chart_png(data)

        from docx import Document
        from docx.shared import Inches
        doc = Document()
        doc.add_heading("Technology Business Management Assessment", 0)
        doc.add_heading("Executive Summary", level=1)

        info = data.get("client_info", {})
        doc.add_paragraph(f"Client: {info.get('name', '')}")
        doc.add_paragraph(f"Industry: {info.get('industry', '')}")
        doc.add_paragraph(f"Date: {info.get('assessment_date', '')}")

        composite = self._weighted_composite(data)
        if composite:
            band = self._get_maturity_band(composite)
            doc.add_heading("Overall Maturity", level=1)
            doc.add_paragraph(f"Score: {composite:.2f} — {band['label']}")

        doc.add_heading("Maturity Profile", level=1)
        doc.add_picture(chart_path, width=Inches(5))

        doc.add_heading("Discipline Scores", level=1)
        for disc in self._enabled_disciplines(data):
            items = []
            for ca in disc.get("capability_areas", []):
                items.extend(ca.get("items", []))
            score = self._score_avg(items)
            doc.add_paragraph(
                f"{disc['name']}: {score:.2f} — {self._get_maturity_band(score)['label']}" if score else f"{disc['name']}: N/A"
            )

        # Top gaps
        target_scores = data.get("target_scores", {})
        gaps = []
        for disc in self._enabled_disciplines(data):
            items = []
            for ca in disc.get("capability_areas", []):
                items.extend(ca.get("items", []))
            current = self._score_avg(items)
            target = target_scores.get(disc["id"], 3.0)
            if current is not None:
                gap = target - current
                if gap > 0:
                    gaps.append({"name": disc["name"], "current": current, "target": target, "gap": gap})
        gaps.sort(key=lambda g: g["gap"], reverse=True)
        if gaps:
            doc.add_heading("Top Priority Gaps", level=1)
            for g in gaps[:5]:
                doc.add_paragraph(f"{g['name']}: Current {g['current']:.2f} → Target {g['target']:.1f} (Gap: {g['gap']:.2f})")

        doc.save(str(output_path))
        return filename

    def export_gap_analysis(self, data: dict) -> str:
        self._ensure_exports_dir()
        filename = f"D-03_TBM_Gap_Analysis_{self._timestamp()}.docx"
        output_path = self.exports_dir / filename
        target_scores = data.get("target_scores", {})

        from docx import Document
        doc = Document()
        doc.add_heading("Technology Business Management Assessment", 0)
        doc.add_heading("Gap Analysis & Roadmap", level=1)

        info = data.get("client_info", {})
        doc.add_paragraph(f"Client: {info.get('name', '')}")
        doc.add_paragraph(f"Date: {info.get('assessment_date', '')}")

        doc.add_heading("Gap Matrix", level=1)
        table = doc.add_table(rows=1, cols=5)
        table.style = "Table Grid"
        headers = ["Discipline", "Current", "Target", "Gap", "Severity"]
        for i, h in enumerate(headers):
            table.rows[0].cells[i].text = h

        for disc in self._enabled_disciplines(data):
            items = []
            for ca in disc.get("capability_areas", []):
                items.extend(ca.get("items", []))
            current = self._score_avg(items)
            target = target_scores.get(disc["id"], 3.0)
            gap = (target - current) if current else None
            severity = "High" if gap and gap > 1.5 else "Medium" if gap and gap > 0.5 else "Low"
            row = table.add_row().cells
            row[0].text = disc["name"]
            row[1].text = f"{current:.2f}" if current else "N/A"
            row[2].text = f"{target:.1f}"
            row[3].text = f"{gap:.2f}" if gap is not None else "N/A"
            row[4].text = severity

        doc.add_heading("Remediation Roadmap", level=1)
        doc.add_heading("30-Day Quick Wins", level=2)
        doc.add_paragraph("Focus on disciplines with scores below 1.75 (Ad Hoc) for immediate improvement opportunities.")
        doc.add_heading("60-Day Improvements", level=2)
        doc.add_paragraph("Establish foundational processes for disciplines scoring 1.75-2.5.")
        doc.add_heading("90-Day Milestones", level=2)
        doc.add_paragraph("Target 'Managed' maturity (2.5+) for critical TBM disciplines.")
        doc.add_heading("6-12 Month Goals", level=2)
        doc.add_paragraph("Achieve 'Optimized' maturity across all disciplines with continuous improvement and automation.")

        doc.save(str(output_path))
        return filename

    def export_workbook(self, data: dict) -> str:
        self._ensure_exports_dir()
        filename = f"D-04_TBM_Workbook_{self._timestamp()}.xlsx"
        output_path = self.exports_dir / filename

        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill

        wb = Workbook()

        # Dashboard sheet
        ws = wb.active
        ws.title = "Dashboard"
        ws["A1"] = "TBM Assessment Dashboard"
        ws["A1"].font = Font(bold=True, size=14)
        ws["A3"] = "Client:"
        ws["B3"] = data.get("client_info", {}).get("name", "")
        ws["A4"] = "Industry:"
        ws["B4"] = data.get("client_info", {}).get("industry", "")
        ws["A5"] = "Assessment Date:"
        ws["B5"] = data.get("client_info", {}).get("assessment_date", "")
        ws["A6"] = "Assessor:"
        ws["B6"] = data.get("client_info", {}).get("assessor", "")

        composite = self._weighted_composite(data)
        ws["A8"] = "Overall Maturity Score:"
        ws["B8"] = round(composite, 2) if composite else "N/A"
        if composite:
            ws["C8"] = self._get_maturity_band(composite)["label"]

        row = 10
        header_fill = PatternFill(start_color="1BA1E2", end_color="1BA1E2", fill_type="solid")
        header_font = Font(bold=True, color="FFFFFF")
        for col, header in enumerate(["Discipline", "Weight", "Score", "Band"], 1):
            cell = ws.cell(row=row, column=col, value=header)
            cell.font = header_font
            cell.fill = header_fill

        for disc in self._enabled_disciplines(data):
            row += 1
            items = []
            for ca in disc.get("capability_areas", []):
                items.extend(ca.get("items", []))
            score = self._score_avg(items)
            ws.cell(row=row, column=1, value=disc["name"])
            ws.cell(row=row, column=2, value=f"{disc.get('weight', 0) * 100:.0f}%")
            ws.cell(row=row, column=3, value=round(score, 2) if score else "N/A")
            if score:
                ws.cell(row=row, column=4, value=self._get_maturity_band(score)["label"])

        # Per-discipline sheets
        for disc in self._enabled_disciplines(data):
            sheet_name = disc["name"][:31]
            ws = wb.create_sheet(title=sheet_name)
            ws["A1"] = disc["name"]
            ws["A1"].font = Font(bold=True, size=12)

            row = 3
            headers = ["Capability Area", "Item #", "Assessment Item", "Score", "Confidence", "Evidence", "Notes"]
            for col, h in enumerate(headers, 1):
                cell = ws.cell(row=row, column=col, value=h)
                cell.font = Font(bold=True)
                cell.fill = PatternFill(start_color="E2E8F0", end_color="E2E8F0", fill_type="solid")

            for ca in disc.get("capability_areas", []):
                for item in ca.get("items", []):
                    row += 1
                    ws.cell(row=row, column=1, value=ca["name"])
                    ws.cell(row=row, column=2, value=item["id"])
                    ws.cell(row=row, column=3, value=item["text"])
                    ws.cell(row=row, column=4, value=item.get("score"))
                    ws.cell(row=row, column=5, value=item.get("confidence", ""))
                    refs = item.get("evidence_references", [])
                    evidence_str = "; ".join(
                        f"{r.get('document', '')} §{r.get('section', '')}"
                        for r in refs if r.get("document")
                    ) if refs else ""
                    ws.cell(row=row, column=6, value=evidence_str)
                    ws.cell(row=row, column=7, value=item.get("notes", ""))

            # Auto-width columns
            for col in ws.columns:
                max_len = 0
                col_letter = col[0].column_letter
                for cell in col:
                    if cell.value:
                        max_len = max(max_len, len(str(cell.value)))
                ws.column_dimensions[col_letter].width = min(max_len + 2, 60)

        wb.save(str(output_path))
        return filename

    def export_outbrief(self, data: dict) -> str:
        self._ensure_exports_dir()
        filename = f"D-05_TBM_Outbrief_{self._timestamp()}.pptx"
        output_path = self.exports_dir / filename
        chart_path = self.generate_radar_chart_png(data)

        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Technology Business Management Assessment"
        subtitle = next((p for p in slide.placeholders if p.placeholder_format.idx == 1), None)
        if subtitle:
            subtitle.text = data.get("client_info", {}).get("name", "")

        # Overview slide
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Assessment Overview"
        composite = self._weighted_composite(data)
        body = next((p for p in slide.placeholders if p.placeholder_format.idx == 1), None)
        if body:
            tf = body.text_frame
            tf.text = f"Client: {data.get('client_info', {}).get('name', '')}"
            tf.add_paragraph().text = f"Industry: {data.get('client_info', {}).get('industry', '')}"
            tf.add_paragraph().text = f"Date: {data.get('client_info', {}).get('assessment_date', '')}"
            tf.add_paragraph().text = f"Assessor: {data.get('client_info', {}).get('assessor', '')}"
            if composite:
                band = self._get_maturity_band(composite)
                tf.add_paragraph().text = f"Overall Score: {composite:.2f} — {band['label']}"

        # Radar chart slide
        blank_idx = min(5, len(prs.slide_layouts) - 1)
        slide = prs.slides.add_slide(prs.slide_layouts[blank_idx])
        slide.shapes.add_picture(chart_path, Inches(3), Inches(0.5), Inches(7), Inches(6.5))

        # Per-discipline slides
        for disc in self._enabled_disciplines(data):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = disc["name"]
            body = next((p for p in slide.placeholders if p.placeholder_format.idx == 1), None)
            if body:
                tf = body.text_frame
                tf.text = ""
                for ca in disc.get("capability_areas", []):
                    items = ca.get("items", [])
                    score = self._score_avg(items)
                    p = tf.add_paragraph()
                    p.text = (
                        f"{ca['name']}: {score:.2f} — {self._get_maturity_band(score)['label']}"
                        if score else f"{ca['name']}: Not scored"
                    )

        prs.save(str(output_path))
        return filename

    def export_heatmap(self, data: dict) -> str:
        self._ensure_exports_dir()
        filename = f"D-06_TBM_Heatmap_{self._timestamp()}.xlsx"
        output_path = self.exports_dir / filename

        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill

        wb = Workbook()
        ws = wb.active
        ws.title = "Heatmap"

        ws["A1"] = "TBM Discipline × Capability Area Heatmap"
        ws["A1"].font = Font(bold=True, size=14)

        score_colors = {
            1: "FF4444", 2: "F97316", 3: "84CC16", 4: "22C55E",
        }

        disciplines = self._enabled_disciplines(data)

        row = 3
        ws.cell(row=row, column=1, value="Discipline").font = Font(bold=True)

        max_cas = 0
        for disc in disciplines:
            max_cas = max(max_cas, len(disc.get("capability_areas", [])))

        for i in range(max_cas):
            ws.cell(row=row, column=i + 2, value=f"CA {i + 1}").font = Font(bold=True)
        ws.cell(row=row, column=max_cas + 2, value="Discipline Avg").font = Font(bold=True)

        for disc in disciplines:
            row += 1
            ws.cell(row=row, column=1, value=disc["name"])
            all_items = []
            for ca_idx, ca in enumerate(disc.get("capability_areas", [])):
                all_items.extend(ca.get("items", []))
                ca_score = self._score_avg(ca.get("items", []))
                cell = ws.cell(row=row, column=ca_idx + 2)
                if ca_score is not None:
                    cell.value = round(ca_score, 2)
                    rounded = min(4, max(1, round(ca_score)))
                    color = score_colors.get(rounded, "FFFFFF")
                    cell.fill = PatternFill(start_color=color, end_color=color, fill_type="solid")
                    if rounded <= 2:
                        cell.font = Font(color="FFFFFF", bold=True)
                    else:
                        cell.font = Font(bold=True)
                else:
                    cell.value = "--"
            domain_avg = self._score_avg(all_items)
            avg_cell = ws.cell(row=row, column=max_cas + 2)
            if domain_avg is not None:
                avg_cell.value = round(domain_avg, 2)
                avg_cell.font = Font(bold=True)

        # Legend
        row += 2
        ws.cell(row=row, column=1, value="Legend:").font = Font(bold=True)
        for score, label in SCORE_LABELS.items():
            row += 1
            cell = ws.cell(row=row, column=1, value=f"{score} — {label}")
            color = score_colors[score]
            cell.fill = PatternFill(start_color=color, end_color=color, fill_type="solid")
            if score <= 2:
                cell.font = Font(color="FFFFFF")

        ws.column_dimensions["A"].width = 30
        for i in range(max_cas + 2):
            col_letter = chr(66 + i) if i < 25 else "Z"
            ws.column_dimensions[col_letter].width = 12

        wb.save(str(output_path))
        return filename

    def export_quick_wins(self, data: dict) -> str:
        self._ensure_exports_dir()
        filename = f"D-07_TBM_Quick_Wins_{self._timestamp()}.docx"
        output_path = self.exports_dir / filename

        from docx import Document
        doc = Document()
        doc.add_heading("Technology Business Management Assessment", 0)
        doc.add_heading("Quick Wins Report", level=1)

        info = data.get("client_info", {})
        doc.add_paragraph(f"Client: {info.get('name', '')}")
        doc.add_paragraph(f"Date: {info.get('assessment_date', '')}")

        doc.add_heading("Methodology", level=1)
        doc.add_paragraph(
            "Quick wins are identified as assessment items scored 1-2 (Ad Hoc or Foundational) "
            "that have the highest potential for improvement. Items are prioritized by discipline weight "
            "multiplied by the score gap to target."
        )

        target_scores = data.get("target_scores", {})
        weights = data.get("scoring_config", {}).get("discipline_weights", {})
        candidates = []
        for disc in self._enabled_disciplines(data):
            weight = weights.get(disc["id"], 0.1)
            target = target_scores.get(disc["id"], 3.0)
            for ca in disc.get("capability_areas", []):
                for item in ca.get("items", []):
                    score = item.get("score")
                    if score is not None and score <= 2 and not item.get("na", False):
                        gap = target - score
                        priority = weight * gap
                        candidates.append({
                            "discipline": disc["name"],
                            "ca": ca["name"],
                            "item_id": item["id"],
                            "text": item["text"],
                            "score": score,
                            "target": target,
                            "gap": gap,
                            "priority": priority,
                            "notes": item.get("notes", ""),
                        })

        candidates.sort(key=lambda c: c["priority"], reverse=True)

        if not candidates:
            doc.add_heading("Results", level=1)
            doc.add_paragraph("No quick win candidates found. All items are scored 3 or above, or are unscored.")
        else:
            from collections import defaultdict
            by_discipline: dict[str, list] = defaultdict(list)
            for c in candidates[:30]:
                by_discipline[c["discipline"]].append(c)

            for disc_name, items in by_discipline.items():
                doc.add_heading(disc_name, level=1)
                table = doc.add_table(rows=1, cols=4)
                table.style = "Table Grid"
                for i, h in enumerate(["Item", "Current", "Target", "Gap"]):
                    table.rows[0].cells[i].text = h
                for item in items:
                    row = table.add_row().cells
                    row[0].text = f"[{item['item_id']}] {item['text'][:80]}"
                    row[1].text = f"{item['score']} ({SCORE_LABELS.get(item['score'], '')})"
                    row[2].text = f"{item['target']:.1f}"
                    row[3].text = f"{item['gap']:.1f}"

        doc.save(str(output_path))
        return filename

    def export_cost_transparency_roadmap(self, data: dict) -> str:
        self._ensure_exports_dir()
        filename = f"D-08_TBM_Cost_Transparency_Roadmap_{self._timestamp()}.docx"
        output_path = self.exports_dir / filename

        from docx import Document
        doc = Document()
        doc.add_heading("Technology Business Management Assessment", 0)
        doc.add_heading("Cost Transparency Roadmap", level=1)

        info = data.get("client_info", {})
        doc.add_paragraph(f"Client: {info.get('name', '')}")
        doc.add_paragraph(f"Industry: {info.get('industry', '')}")
        doc.add_paragraph(f"Date: {info.get('assessment_date', '')}")

        composite = self._weighted_composite(data)

        # Current maturity assessment of cost transparency practices
        doc.add_heading("Current Cost Transparency Maturity", level=1)
        if composite:
            band = self._get_maturity_band(composite)
            doc.add_paragraph(f"Overall TBM Maturity: {composite:.2f} — {band['label']}")
        else:
            doc.add_paragraph("Overall TBM Maturity: Not yet assessed")

        # Find cost-related disciplines
        cost_disciplines = []
        for disc in self._enabled_disciplines(data):
            items = []
            for ca in disc.get("capability_areas", []):
                items.extend(ca.get("items", []))
            score = self._score_avg(items)
            cost_disciplines.append({"name": disc["name"], "score": score, "id": disc["id"]})

        doc.add_paragraph("Discipline-level maturity for cost transparency:")
        for cd in cost_disciplines:
            if cd["score"] is not None:
                band = self._get_maturity_band(cd["score"])
                doc.add_paragraph(f"  • {cd['name']}: {cd['score']:.2f} — {band['label']}", style="List Bullet")
            else:
                doc.add_paragraph(f"  • {cd['name']}: Not scored", style="List Bullet")

        # Phased roadmap for TBM taxonomy adoption
        doc.add_heading("Phased Roadmap for TBM Taxonomy Adoption", level=1)

        doc.add_heading("Phase 1: Foundation (0-3 months)", level=2)
        doc.add_paragraph("• Establish IT cost categorization using TBM taxonomy (Cost Pools, IT Towers, Services)")
        doc.add_paragraph("• Identify and catalog all IT cost sources")
        doc.add_paragraph("• Map existing chart of accounts to TBM cost pools")
        doc.add_paragraph("• Define initial allocation rules for shared costs")

        doc.add_heading("Phase 2: Standardization (3-6 months)", level=2)
        doc.add_paragraph("• Implement TBM taxonomy across financial systems")
        doc.add_paragraph("• Establish automated data collection from ITSM, ITAM, and financial systems")
        doc.add_paragraph("• Create standard cost models for IT towers and services")
        doc.add_paragraph("• Begin regular cost transparency reporting to IT leadership")

        doc.add_heading("Phase 3: Optimization (6-12 months)", level=2)
        doc.add_paragraph("• Implement unit cost metrics for key IT services")
        doc.add_paragraph("• Establish benchmarking against industry peers")
        doc.add_paragraph("• Create business-facing IT cost views tied to business capabilities")
        doc.add_paragraph("• Implement showback/chargeback mechanisms")

        doc.add_heading("Phase 4: Value Realization (12+ months)", level=2)
        doc.add_paragraph("• Enable data-driven investment decisions using TBM insights")
        doc.add_paragraph("• Optimize IT spending through continuous cost analysis")
        doc.add_paragraph("• Align IT investments with business value delivery")
        doc.add_paragraph("• Establish mature total cost of ownership (TCO) models")

        # Cost allocation maturity progression recommendations
        doc.add_heading("Cost Allocation Maturity Progression", level=1)

        doc.add_heading("Ad Hoc → Foundational", level=2)
        doc.add_paragraph("• Move from spreadsheet-based tracking to centralized cost repository")
        doc.add_paragraph("• Establish basic cost pool categorization")
        doc.add_paragraph("• Define ownership for cost data quality")

        doc.add_heading("Foundational → Managed", level=2)
        doc.add_paragraph("• Implement systematic allocation methodologies")
        doc.add_paragraph("• Automate data collection from primary source systems")
        doc.add_paragraph("• Establish regular cost review cadence with stakeholders")

        doc.add_heading("Managed → Optimized", level=2)
        doc.add_paragraph("• Achieve full automation of cost allocation processes")
        doc.add_paragraph("• Implement predictive cost modeling and forecasting")
        doc.add_paragraph("• Drive continuous improvement through analytics and benchmarking")

        # Quick wins in cost transparency
        doc.add_heading("Quick Wins in Cost Transparency", level=1)

        target_scores = data.get("target_scores", {})
        quick_wins_found = False
        for disc in self._enabled_disciplines(data):
            target = target_scores.get(disc["id"], 3.0)
            for ca in disc.get("capability_areas", []):
                for item in ca.get("items", []):
                    score = item.get("score")
                    if score is not None and score <= 2 and not item.get("na", False):
                        if not quick_wins_found:
                            quick_wins_found = True
                        doc.add_paragraph(
                            f"• [{disc['name']} / {ca['name']}] {item['text'][:100]} "
                            f"(Current: {score}, Target: {target:.0f})",
                            style="List Bullet"
                        )

        if not quick_wins_found:
            doc.add_paragraph("No immediate quick wins identified. All scored items are at Managed level or above.")

        doc.save(str(output_path))
        return filename

    def export_all(self, data: dict) -> list[str]:
        return [
            self.export_findings(data),
            self.export_executive_summary(data),
            self.export_gap_analysis(data),
            self.export_workbook(data),
            self.export_outbrief(data),
            self.export_heatmap(data),
            self.export_quick_wins(data),
            self.export_cost_transparency_roadmap(data),
        ]
